from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import requests
from io import BytesIO
import re
import json

class SlideGenerator:
    def __init__(self):
        pass

    def download_image(self, query):
        if not query:
            return None
            
        print(f"Searching for image (Bing): {query}")
        try:
            # Bing Image Search
            url = f"https://www.bing.com/images/search?q={query}&first=1"
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
            }
            
            response = requests.get(url, headers=headers, timeout=10)
            if response.status_code != 200:
                print(f"Bing Search failed: {response.status_code}")
                return None
                
            # Extract image URL using regex
            # Look for murl":"https://..."
            # The attribute is often HTML encoded, but in the script tag or data attribute it might be plain.
            # Common pattern in Bing: murl":"(URL)"
            
            matches = re.search(r'murl":"(https?://[^"]+)"', response.text)
            if not matches:
                # Try alternative pattern (HTML encoded)
                matches = re.search(r'murl&quot;:&quot;(https?://[^&]+)&quot;', response.text)
                
            if not matches:
                print(f"No images found for query: {query}")
                return None
                
            image_url = matches.group(1)
            # Decode if it was HTML encoded (simple check)
            if "&quot;" in image_url:
                image_url = image_url.replace("&quot;", "")
                
            print(f"Found image URL: {image_url}")
            
            # Download image
            img_response = requests.get(image_url, headers=headers, timeout=10)
            if img_response.status_code == 200:
                return BytesIO(img_response.content)
            else:
                print(f"Failed to download image. Status: {img_response.status_code}")
                
        except Exception as e:
            print(f"Error getting image for query '{query}': {e}")
            
        return None

    def create_presentation(self, data: dict, output_file: str = "lecture_slides.pptx") -> str:
        prs = Presentation()

        # --- Title Slide ---
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = data.get("topic", "University Lecture")
        
        # Construct subtitle with metadata
        uni = data.get("university", "")
        subj = data.get("subject", "")
        lecturer = data.get("lecturer", "")
        
        subtitle_text = []
        if uni: subtitle_text.append(uni)
        if subj: subtitle_text.append(subj)
        if lecturer: subtitle_text.append(f"Lecturer: {lecturer}")
        
        subtitle.text = "\n".join(subtitle_text)

        # --- Content Slides ---
        for slide_data in data.get("slides", []):
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = slide_data.get("title", "Untitled Slide")
            
            # Text Frame
            tf = body_shape.text_frame
            tf.text = ""  # clear existing text

            # Handle Image
            image_stream = None
            # Check for both keys to support backward compatibility or manual edits
            query = slide_data.get("image_search_query") or slide_data.get("image_prompt")
            
            if query:
                image_stream = self.download_image(query)

            if image_stream:
                # Resize text box to 50% width to make room for image
                body_shape.width = Inches(5.0)
                body_shape.top = Inches(2.0)
                body_shape.left = Inches(0.5)
                
                # Add Image
                try:
                    # Position image on the right side
                    slide.shapes.add_picture(image_stream, left=Inches(5.8), top=Inches(2.0), width=Inches(4.0))
                except Exception as e:
                    print(f"Error adding picture to slide: {e}")
                    # If image fails, revert text box to full width
                    body_shape.width = Inches(9.0)
            else:
                # Ensure full width if no image
                body_shape.width = Inches(9.0)
                body_shape.top = Inches(2.0)
                body_shape.left = Inches(0.5)

            # Add Content
            tf.word_wrap = True
            
            for point in slide_data.get("content", []):
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                
                # Adjust font size based on layout
                if image_stream:
                    p.font.size = Pt(18)  # Smaller font if sharing space with image
                else:
                    p.font.size = Pt(24)  # Standard font size

            # Add speaker notes
            if "speaker_notes" in slide_data:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = slide_data["speaker_notes"]

        # Save
        os.makedirs(os.path.dirname(os.path.abspath(output_file)), exist_ok=True)
        prs.save(output_file)
        return output_file
